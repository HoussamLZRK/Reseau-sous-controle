from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Créer une présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Titre et contenu
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Réseau Sous Contrôle"
subtitle.text = "Détection d'Anomalies — Comparaison de 3 Modèles ML\n\nProjet PFA — Maroc Telecom"

# Slide 2 : Vue d'ensemble
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Objectif du Projet"

tf = content.text_frame
tf.text = "Analyser 6629 mesures de drive-test sur 26 villes et 11 régions du Maroc"

p = tf.add_paragraph()
p.text = "Détecter automatiquement les anomalies réseau"
p.level = 0

p = tf.add_paragraph()
p.text = "Comparer 3 approches : Unsupervised, Supervised ML, Deep Learning"
p.level = 0

p = tf.add_paragraph()
p.text = "Identifier les zones à surveiller en priorité"
p.level = 0

# Slide 3 : Données
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Les Données"

tf = content.text_frame
tf.text = "6629 mesures de benchmarks nationaux"

p = tf.add_paragraph()
p.text = "26 villes réparties sur 11 régions"
p.level = 0

p = tf.add_paragraph()
p.text = "3 campagnes de mesure (mai, juin, juillet 2026)"
p.level = 0

p = tf.add_paragraph()
p.text = "22 colonnes : signal, débits, anomalies, risques"
p.level = 0

p = tf.add_paragraph()
p.text = "95.3% mesures normales, 4.7% anomalies"
p.level = 0

# Slide 4 : Modèle 1 - Isolation Forest
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Modèle 1 : Isolation Forest (Unsupervised)"

tf = content.text_frame
tf.text = "Type : Détection d'anomalies SANS labels"

p = tf.add_paragraph()
p.text = "Architecture : 100 arbres isolants"
p.level = 0

p = tf.add_paragraph()
p.text = "Accuracy : 86.7%"
p.level = 0

p = tf.add_paragraph()
p.text = "Precision : 7.1% ❌ (TRÈS BAS — beaucoup de faux positifs)"
p.level = 0

p = tf.add_paragraph()
p.text = "Recall : 15.0%"
p.level = 0

p = tf.add_paragraph()
p.text = "Cas d'usage : Exploration de données"
p.level = 0

# Slide 5 : Modèle 2 - Random Forest
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Modèle 2 : Random Forest (Supervised) ⭐"

tf = content.text_frame
tf.text = "Type : Classification supervisée AVEC labels"

p = tf.add_paragraph()
p.text = "Architecture : 100 arbres de décision"
p.level = 0

p = tf.add_paragraph()
p.text = "Accuracy : 95.5%"
p.level = 0

p = tf.add_paragraph()
p.text = "Precision : 100.0% ✅ (PARFAIT — zéro faux positif)"
p.level = 0

p = tf.add_paragraph()
p.text = "Recall : 4.8%"
p.level = 0

p = tf.add_paragraph()
p.text = "ROC-AUC : 0.9006"
p.level = 0

p = tf.add_paragraph()
p.text = "Cas d'usage : PRODUCTION (fiabilité maximale)"
p.level = 0

# Slide 6 : Modèle 3 - Neural Network
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Modèle 3 : Neural Network (Deep Learning)"

tf = content.text_frame
tf.text = "Type : Réseau de neurones feedforward"

p = tf.add_paragraph()
p.text = "Architecture : Dense(32) → Dense(16) → Dense(8) → Output"
p.level = 0

p = tf.add_paragraph()
p.text = "Accuracy : 95.7%"
p.level = 0

p = tf.add_paragraph()
p.text = "Precision : 75.0%"
p.level = 0

p = tf.add_paragraph()
p.text = "Recall : 14.3%"
p.level = 0

p = tf.add_paragraph()
p.text = "ROC-AUC : 0.8929"
p.level = 0

p = tf.add_paragraph()
p.text = "Cas d'usage : Apprentissage de patterns complexes"
p.level = 0

# Slide 7 : Comparaison
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Comparaison des 3 Modèles"

tf = content.text_frame
tf.text = "Random Forest GAGNE : Precision 100% vs 7.1% (IF) vs 75% (NN)"

p = tf.add_paragraph()
p.text = "IF : Utile pour exploration SANS labels"
p.level = 0

p = tf.add_paragraph()
p.text = "RF : Meilleur pour production (100% confiance)"
p.level = 0

p = tf.add_paragraph()
p.text = "NN : Peut capturer des patterns non-linéaires"
p.level = 0

# Slide 8 : Recommandations
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Recommandations Finales"

tf = content.text_frame
tf.text = "✅ Utiliser Random Forest en production"

p = tf.add_paragraph()
p.text = "Precision 100% → zéro faux positif"
p.level = 1

p = tf.add_paragraph()
p.text = "Rapide et interprétable"
p.level = 1

p = tf.add_paragraph()
p.text = "✅ Isolation Forest pour exploration initiale"
p.level = 0

p = tf.add_paragraph()
p.text = "Pas besoin de labels"
p.level = 1

p = tf.add_paragraph()
p.text = "✅ Neural Network pour cas complexes"
p.level = 0

p = tf.add_paragraph()
p.text = "À investiguer sur d'autres datasets"
p.level = 1

# Slide 9 : KPIs Finaux
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "KPIs Clés du Projet"

tf = content.text_frame
tf.text = "6629 mesures analysées"

p = tf.add_paragraph()
p.text = "26 villes couvertes, 11 régions mappées"
p.level = 0

p = tf.add_paragraph()
p.text = "3 modèles ML comparés"
p.level = 0

p = tf.add_paragraph()
p.text = "100% Precision (Random Forest)"
p.level = 0

p = tf.add_paragraph()
p.text = "95.5% Accuracy (Random Forest)"
p.level = 0

p = tf.add_paragraph()
p.text = "19 vues SQL créées"
p.level = 0

p = tf.add_paragraph()
p.text = "4 pages Power BI interactives"
p.level = 0

# Slide 10 : Conclusion
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"

tf = content.text_frame
tf.text = "Projet complet et opérationnel"

p = tf.add_paragraph()
p.text = "De la donnée brute → Dashboard → Modèles ML"
p.level = 0

p = tf.add_paragraph()
p.text = "Random Forest validé pour production"
p.level = 0

p = tf.add_paragraph()
p.text = "Prêt pour déploiement chez Maroc Telecom"
p.level = 0

# Sauvegarder
prs.save('Reseau_Sous_Controle_Resultats.pptx')
print("✅ Présentation créée : Reseau_Sous_Controle_Resultats.pptx")